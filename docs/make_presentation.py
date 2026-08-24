"""
Build the APOLLO-M mid-evaluation presentation.

Built on the team's own proposal deck so the university master, background and
Department of Computer Sciences logo carry over unchanged: the proposal's
content slides are removed and new ones are added on the same branded layout,
rather than the branding being rebuilt from scratch.

The template's usable area is narrow — the background art puts a lavender band
across the bottom 1.1 inches and the logo sits inside it — so every slide keeps
its content between y=1.15 and y=6.15 and nothing is drawn over the logo.

Wording is deliberately plain and short: lines meant to be explained aloud from
memory, not read off the slide. Every figure quoted is measured and traceable to
a script in this repository.

    python docs/make_presentation.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

HERE = Path(__file__).resolve().parent
IMG = HERE / "slides_img"
FIG = HERE / "figures"
SRC = Path(r"C:\Users\kasha\Downloads\FINALYEARPROJECTPPT (3).pptx")
OUT = HERE / "APOLLO-M_Mid_Evaluation_Presentation.pptx"

FONT = "Times New Roman"

# Pastel palette, chosen to sit with the template's lavender band.
# No dark green and no red anywhere — apricot stands in wherever a warning
# colour would normally be used.
INK = RGBColor(0x33, 0x41, 0x5C)        # deep slate, body text
TITLE_C = RGBColor(0x2A, 0x3A, 0x5C)
MUTED = RGBColor(0x6B, 0x7A, 0x99)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLUE = RGBColor(0x7F, 0xA8, 0xD4)
LAV = RGBColor(0xA9, 0x9B, 0xD1)
TEAL = RGBColor(0x7F, 0xC0, 0xBA)
APRICOT = RGBColor(0xE8, 0xAF, 0x87)
SAGE = RGBColor(0xA9, 0xC7, 0x9E)
CARD_BG = RGBColor(0xF5, 0xF8, 0xFC)

# Content band. Above: nothing. Below: the template's band and the logo.
TOP = 1.18
BOT = 6.15

prs = Presentation(str(SRC))

# Drop the proposal's slides but keep masters and layouts, where the branding
# actually lives. The relationship has to be dropped as well as the id —
# removing only the id leaves the old parts in the package, which then collide
# with the new slide1.xml, slide2.xml ... on save.
_sld_lst = prs.slides._sldIdLst
for sld_id in list(_sld_lst):
    prs.part.drop_rel(sld_id.rId)
    _sld_lst.remove(sld_id)

# The branded layout is the one carrying the background art and the logo.
BRANDED = next(
    (lay for lay in prs.slide_layouts
     if sum(1 for sh in lay.shapes if sh.shape_type == 13) >= 2),
    prs.slide_layouts[1],
)


def slide(title=None):
    """A slide on the branded layout, with the template's own placeholders
    removed so nothing shows 'Click to add text' and our own boxes control
    every position."""
    s = prs.slides.add_slide(BRANDED)
    for ph in list(s.placeholders):
        ph._element.getparent().remove(ph._element)
    if title:
        tb = s.shapes.add_textbox(Inches(0.55), Inches(0.36), Inches(8.9), Inches(0.7))
        p = tb.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = title
        r.font.name, r.font.size, r.font.bold = FONT, Pt(30), True
        r.font.color.rgb = TITLE_C
        rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.58), Inches(1.0),
                                  Inches(1.5), Inches(0.05))
        rule.fill.solid()
        rule.fill.fore_color.rgb = LAV
        rule.line.fill.background()
        rule.shadow.inherit = False
    return s


def txt(s, x, y, w, h, lines, size=18, color=INK, bold=False, italic=False,
        align=PP_ALIGN.LEFT, bullet=False, space=10):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    items = lines if isinstance(lines, list) else [lines]
    for i, line in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        r = p.add_run()
        r.text = ("\u2022   " + line) if bullet else line
        r.font.name, r.font.size = FONT, Pt(size)
        r.font.bold, r.font.italic = bold, italic
        r.font.color.rgb = color
    return tb


def card(s, x, y, w, h, big, label, accent=BLUE, big_pt=30):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                            Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = CARD_BG
    sh.line.color.rgb = accent
    sh.line.width = Pt(1.75)
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = big
    r.font.name, r.font.size, r.font.bold = FONT, Pt(big_pt), True
    r.font.color.rgb = TITLE_C
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = label
    r2.font.name, r2.font.size = FONT, Pt(12)
    r2.font.color.rgb = MUTED
    return sh


def pill(s, x, y, w, h, text, fill, size=14):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                            Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name, r.font.size, r.font.bold = FONT, Pt(size), True
    r.font.color.rgb = WHITE
    return sh


def pic(s, path, x, y, w=None, h=None):
    p = Path(path)
    if not p.exists():
        raise FileNotFoundError(p)
    kw = {}
    if w:
        kw["width"] = Inches(w)
    if h:
        kw["height"] = Inches(h)
    return s.shapes.add_picture(str(p), Inches(x), Inches(y), **kw)


# ───────────────────────────── 1 · Title ─────────────────────────────────
s = slide()
txt(s, 0.6, 1.55, 8.8, 1.1, ["APOLLO-M"], size=50, bold=True,
    align=PP_ALIGN.CENTER, color=TITLE_C)
txt(s, 0.9, 2.62, 8.2, 0.9,
    ["An AI Framework for Forecasting Instability in Digital Communities"],
    size=19, align=PP_ALIGN.CENTER, color=INK)
rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.0), Inches(3.55),
                          Inches(2.0), Inches(0.05))
rule.fill.solid()
rule.fill.fore_color.rgb = LAV
rule.line.fill.background()
rule.shadow.inherit = False
txt(s, 0.6, 3.85, 8.8, 1.5,
    ["Kashaf Fatima  (84381)      Rizwan Saleem  (84010)      Aaqib Mehmood  (82081)",
     "Supervisor:  Mr. Saghir Ahmed",
     "Mid-Year Evaluation  \u00b7  Final Year Project"],
    size=15, align=PP_ALIGN.CENTER, color=MUTED, space=8)
txt(s, 0.6, 5.45, 8.8, 0.5,
    ["Live:  apollo-m.streamlit.app"],
    size=14, align=PP_ALIGN.CENTER, color=BLUE, bold=True)

# ───────────────────────────── 2 · Overview ──────────────────────────────
s = slide("What We Will Cover")
left = ["The problem we are solving", "Our idea", "What we built",
        "Our key finding", "Does it actually predict?"]
right = ["Running on real Reddit data", "Who would use it, and why",
         "Industry and the SDGs", "Timeline and limitations", "Next steps"]
y = TOP + 0.15
for i, (a, b) in enumerate(zip(left, right)):
    pill(s, 0.55, y, 0.42, 0.42, str(i + 1), BLUE, size=13)
    txt(s, 1.1, y - 0.02, 3.6, 0.5, [a], size=16.5)
    pill(s, 5.05, y, 0.42, 0.42, str(i + 6), LAV, size=13)
    txt(s, 5.6, y - 0.02, 3.9, 0.5, [b], size=16.5)
    y += 0.9

# ───────────────────────────── 3 · Problem ───────────────────────────────
s = slide("The Problem")
txt(s, 0.55, TOP, 8.9, 0.9,
    ["Online communities can break down. Hostility rises, the people worth "
     "keeping leave, and moderators are left firefighting."], size=19)
txt(s, 0.55, TOP + 1.05, 8.9, 2.4,
    ["Today's tools judge one comment at a time, after it is posted.",
     "They describe how a community looks now \u2014 not where it is heading.",
     "By the time a dashboard turns red, the harm has already happened."],
    size=18, bullet=True, space=16)
card(s, 0.55, 4.85, 8.9, 1.2,
     "Moderators are given volume, not priority.",
     "the gap we set out to close", LAV, big_pt=22)

# ───────────────────────────── 4 · The idea ──────────────────────────────
s = slide("Our Idea \u2014 a Lesson from Universe 25")
txt(s, 0.55, TOP, 8.9, 1.6,
    ["From 1968 to 1973, John Calhoun gave a mouse colony unlimited food, water "
     "and safety \u2014 but limited space. The population grew, then collapsed. "
     "The collapse was described in detail afterwards. Nothing measured it "
     "coming."], size=18)
card(s, 0.55, 2.95, 8.9, 1.1,
     "Online communities have the same problem \u2014 but they are fully recorded.",
     "every comment, every author, every timestamp", TEAL, big_pt=19)
txt(s, 0.55, 4.3, 8.9, 1.7,
    ["So the early warning signs are already in the data. Nobody was looking "
     "for them.",
     "We treat instability as something to forecast, not something to classify."],
    size=18, bullet=True, space=14)

# ───────────────────────────── 5 · What we built ─────────────────────────
s = slide("What We Built")
pic(s, FIG / "apollo_architecture.png", 0.6, TOP, h=4.75)
txt(s, 6.3, TOP + 0.15, 3.15, 4.4,
    ["Five layers.",
     "Each one is measured on its own, so we can say which part works and "
     "which does not.",
     "Deterministic code and trained models make every decision. The language "
     "model only writes the explanation \u2014 it never changes a score."],
    size=15, space=14)

# ───────────────────────────── 6 · How it works ──────────────────────────
s = slide("How It Works \u2014 in Plain Terms")
rows = [
    ("MICRO", "Score every comment for toxicity", "toxic-bert", BLUE),
    ("MESO", "Combine them into a health score per community",
     "health index + instability", TEAL),
    ("PATTERN", "Group similar communities without labels",
     "K-Means, DBSCAN", LAV),
    ("MACRO", "Forecast the next five days, with confidence bands",
     "Temporal Fusion Transformer", SAGE),
    ("ACT", "Raise a graded alert and a recommended action",
     "four alert bands", APRICOT),
]
y = TOP - 0.02
for tag, what, how, col in rows:
    pill(s, 0.55, y, 1.45, 0.8, tag, col, size=13)
    txt(s, 2.2, y + 0.03, 4.8, 0.75, [what], size=16.5)
    txt(s, 7.05, y + 0.08, 2.45, 0.7, [how], size=12, color=MUTED, italic=True)
    y += 0.97

# ───────────────────────────── 7 · The finding ───────────────────────────
s = slide("Our Key Finding")
txt(s, 0.55, TOP - 0.08, 8.9, 0.5,
    ["We tested our own health score \u2014 it did worse than an ingredient "
     "inside it."], size=17.5, bold=True)
pic(s, IMG / "finding_chart.png", 0.9, TOP + 0.42, w=8.2)
txt(s, 0.55, 5.82, 8.9, 0.5,
    ["A health score says how bad a community is now. Instability is how fast "
     "it is changing."],
    size=14.5, color=MUTED, italic=True, align=PP_ALIGN.CENTER)

# ───────────────────────────── 8 · Confirmed twice ───────────────────────
s = slide("The Same Failure, Found Twice")
txt(s, 0.55, TOP, 8.9, 0.85,
    ["Our graph model was built separately, with different inputs. It scored "
     "below chance on the very same task."], size=18)
card(s, 0.55, TOP + 0.95, 4.3, 1.45, "0.575", "Community Health Index", APRICOT)
card(s, 5.15, TOP + 0.95, 4.3, 1.45, "0.415", "GraphSAGE structural risk", APRICOT)
txt(s, 0.55, TOP + 2.65, 8.9, 2.1,
    ["Both were trained on how a community looks now, not on how it is changing.",
     "One weak score could be a mistake in one formula.",
     "Two, built independently, show that the whole type of measure is wrong."],
    size=17.5, bullet=True, space=14)

# ───────────────────────────── 9 · Forecasting ───────────────────────────
s = slide("Does It Actually Predict?")
txt(s, 0.55, TOP - 0.05, 8.9, 0.6,
    ["We hid 15 destabilising communities among 60. The forecaster never sees "
     "the answer key."], size=17.5)
card(s, 0.55, TOP + 0.65, 2.85, 1.4, "15 / 15", "found correctly", TEAL)
card(s, 3.62, TOP + 0.65, 2.8, 1.4, "0.732", "trend ranking (ROC-AUC)", BLUE)
card(s, 6.64, TOP + 0.65, 2.8, 1.4, "5 days", "forecast horizon", LAV)
txt(s, 0.55, TOP + 2.3, 8.9, 2.4,
    ["The predicted direction was right for every planted community.",
     "Predicted daily trend:  getting worse  +0.0024   \u00b7   stable  +0.0009   "
     "\u00b7   improving  \u22120.0008",
     "Missing a community that is deteriorating costs far more than checking a "
     "quiet one, so we tune for recall."],
    size=17, bullet=True, space=14)

# ───────────────────────────── 10 · Real data ────────────────────────────
s = slide("It Runs on Real Reddit Data")
# The counts run down the left so the dashboard screenshot keeps enough width
# to stay legible without reaching the logo band.
_cy = TOP - 0.02
for _big, _lab, _col in [("116,854", "real comments", BLUE),
                         ("24", "communities", TEAL),
                         ("64,416", "real authors", LAV),
                         ("40", "days of history", SAGE)]:
    card(s, 0.55, _cy, 2.25, 1.1, _big, _lab, _col, big_pt=22)
    _cy += 1.22
pic(s, IMG / "image24.png", 3.15, TOP + 0.52, w=6.3)
txt(s, 3.15, 5.5, 6.3, 0.65,
    ["Collected through Arctic Shift \u2014 no API key needed. r/PublicFreakout "
     "came out least healthy, and we never told it to."],
    size=12.5, color=MUTED, italic=True, align=PP_ALIGN.CENTER, space=4)

# ───────────────────────────── 11 · Live system ──────────────────────────
s = slide("The Live System")
# Four surfaces, four screenshots. Each is scaled to the same height and
# centred in its half so the grid lines up whatever the source aspect ratio.
_shots = [
    ("image34.png", "REST API with interactive docs", 0.55, TOP - 0.05),
    ("image25.png", "Grafana monitoring, live metrics", 5.05, TOP - 0.05),
    ("image23.png", "Live Reddit fetch, on demand", 0.55, TOP + 2.35),
    ("image46.png", "Five-day forecast with uncertainty", 5.05, TOP + 2.35),
]
for _name, _cap, _x, _y in _shots:
    _p = pic(s, IMG / _name, _x, _y, h=1.85)
    _p.left = Inches(_x + (4.4 - _p.width / 914400) / 2)
    txt(s, _x, _y + 1.92, 4.4, 0.35, [_cap], size=12, color=MUTED,
        italic=True, align=PP_ALIGN.CENTER)
txt(s, 0.55, 5.78, 8.9, 0.5,
    ["apollo-m.streamlit.app    ·    apollo-api-tllm.onrender.com/docs"
     "    ·    cerebro-sandy-beta.vercel.app"],
    size=13.5, color=BLUE, bold=True, align=PP_ALIGN.CENTER)

# ───────────────────────────── 12 · Evidence ─────────────────────────────
s = slide("Evidence That It Works")
pic(s, IMG / "image8.png", 0.5, TOP + 0.05, w=9.0)
txt(s, 0.55, TOP + 3.05, 8.9, 2.0,
    ["14 automated tests on the forecasting pipeline; 67 on the companion "
     "system, CEREBRO.",
     "Toxicity classifier: 90.1% accuracy on real labelled data held back from "
     "training.",
     "Every number in our report comes from a script that can be run again."],
    size=16.5, bullet=True, space=13)

# ───────────────────────────── 13 · Buyers ───────────────────────────────
s = slide("Who Would Use It \u2014 and Why")
groups = [
    ("Trust and Safety teams", "Reddit, Discord, Twitch, forum hosts",
     "They cannot watch 500 communities at once. We tell them which ten to "
     "watch this week.", BLUE),
    ("Community platform vendors", "Discourse, Circle, Mighty Networks",
     "A health feature they can resell to their own customers.", TEAL),
    ("Brand safety and advertising", "ad-verification firms",
     "Advertisers pay to avoid deteriorating spaces \u2014 before the story "
     "breaks, not after.", LAV),
    ("Large organisations", "internal Slack and Teams communities",
     "Culture problems usually surface far too late.", SAGE),
]
y = TOP - 0.05
for name, who, why, col in groups:
    bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55),
                             Inches(y), Inches(0.13), Inches(1.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = col
    bar.line.fill.background()
    bar.shadow.inherit = False
    txt(s, 0.85, y - 0.06, 3.4, 0.45, [name], size=16, bold=True)
    txt(s, 0.85, y + 0.35, 3.4, 0.5, [who], size=11.5, color=MUTED, italic=True)
    txt(s, 4.45, y, 5.0, 1.0, [why], size=14)
    y += 1.17
txt(s, 0.55, 5.72, 8.9, 0.5,
    ["Why buy rather than build: the hard part is knowing what to measure "
     "\u2014 and we measured it."],
    size=14, color=MUTED, italic=True, align=PP_ALIGN.CENTER)

# ───────────────────────────── 14 · Industry + SDG ───────────────────────
s = slide("Industry Relevance and the SDGs")
txt(s, 0.55, TOP - 0.05, 8.9, 0.55,
    ["Industry:  Trust and Safety \u2014 where social platforms, brand safety "
     "and regulation meet."], size=17, bold=True)
card(s, 0.55, TOP + 0.6, 4.35, 1.5, "SDG 16",
     "Peace, Justice and Strong Institutions", TEAL)
card(s, 5.1, TOP + 0.6, 4.35, 1.5, "SDG 9",
     "Industry, Innovation and Infrastructure", BLUE)
txt(s, 0.6, TOP + 2.3, 4.25, 2.0,
    ["Target 16.1 \u2014 reduce violence.",
     "Online hostility and harassment come before real-world harm. Acting "
     "earlier means acting before it spreads."], size=14.5, space=10)
txt(s, 5.15, TOP + 2.3, 4.25, 2.0,
    ["Target 9.5 \u2014 build research capability.",
     "The API, the monitoring stack and the published benchmark are reusable "
     "infrastructure, not a one-off model."], size=14.5, space=10)

# ───────────────────────────── 15 · Timeline ─────────────────────────────
s = slide("Timeline \u2014 Planned Against Delivered")
pic(s, IMG / "timeline_chart.png", 0.75, TOP - 0.08, w=8.5)
txt(s, 0.55, 5.78, 8.9, 0.5,
    ["Weeks 4 to 6 were not in the plan \u2014 we found real problems in our own "
     "data and stopped to fix them."],
    size=14, color=MUTED, italic=True, align=PP_ALIGN.CENTER)

# ───────────────────────────── 16 · Limitations ──────────────────────────
s = slide("What We Are Honest About")
txt(s, 0.55, TOP - 0.02, 8.9, 4.5,
    ["Our ground truth is planted, not observed \u2014 we know which "
     "communities we made unstable, because nobody publishes a list of real "
     "ones.",
     "So the instability score of 1.000 shows the pipeline is wired correctly, "
     "not that it is that accurate in the real world.",
     "The misinformation model was trained on news articles. On short posts it "
     "flags almost everything, so we publish no rate from it.",
     "Evaluation is single-fold; we do not yet report confidence intervals.",
     "English only."],
    size=16, bullet=True, space=15)
txt(s, 0.55, 5.72, 8.9, 0.5,
    ["Stating the limits is what makes the other numbers believable."],
    size=14.5, color=MUTED, italic=True, align=PP_ALIGN.CENTER)

# ───────────────────────────── 17 · Next steps ───────────────────────────
s = slide("Next Steps")
items = [
    ("Real ground truth",
     "Reddit publishes quarantine and ban dates. Train up to 60 days before "
     "each one, then check whether we flagged it first."),
    ("Longer history",
     "Extend to 18\u201324 months per community using bulk archives."),
    ("Wire the graph model into the health score",
     "Now that it is measured, we know what it does and does not add."),
    ("Retrain on the right data",
     "Misinformation on social text; the recommender on real moderator "
     "decisions."),
    ("Publish the benchmark",
     "Release the corpus and ground truth so others can reproduce our results."),
]
y = TOP - 0.05
for i, (head, body) in enumerate(items, 1):
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.55), Inches(y),
                           Inches(0.48), Inches(0.48))
    c.fill.solid()
    c.fill.fore_color.rgb = LAV
    c.line.fill.background()
    c.shadow.inherit = False
    tf = c.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = str(i)
    r.font.name, r.font.size, r.font.bold = FONT, Pt(14), True
    r.font.color.rgb = WHITE
    txt(s, 1.2, y - 0.08, 8.25, 0.45, [head], size=16, bold=True)
    txt(s, 1.2, y + 0.32, 8.25, 0.6, [body], size=13, color=MUTED)
    y += 0.99

# ───────────────────────────── 18 · Conclusion ───────────────────────────
s = slide("Conclusion")
txt(s, 0.55, TOP - 0.05, 8.9, 0.55,
    ["We set out to predict community breakdown instead of reacting to it."],
    size=18.5, bold=True)
txt(s, 0.55, TOP + 0.65, 8.9, 3.3,
    ["We built it \u2014 five layers, from a single comment to a five-day "
     "forecast, deployed and publicly reachable.",
     "We tested it honestly, and found our own health score was the weakest "
     "signal we had.",
     "We fixed it \u2014 then a second model failed the same way, which turned "
     "one result into a general one.",
     "It runs on real Reddit data, collected without any API key."],
    size=17, bullet=True, space=14)
card(s, 0.55, 5.0, 8.9, 1.1,
     "Anyone can build a pipeline. We measured ours and reported what it said.",
     "that is our contribution", TEAL, big_pt=18)

# ───────────────────────────── 19 · Thank you ────────────────────────────
s = slide()
txt(s, 0.6, 1.9, 8.8, 0.95, ["Thank You"], size=44, bold=True,
    align=PP_ALIGN.CENTER, color=TITLE_C)
txt(s, 0.6, 2.9, 8.8, 0.5, ["Questions?"], size=21,
    align=PP_ALIGN.CENTER, color=MUTED)
rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.25), Inches(3.65),
                          Inches(1.5), Inches(0.05))
rule.fill.solid()
rule.fill.fore_color.rgb = LAV
rule.line.fill.background()
rule.shadow.inherit = False
txt(s, 0.6, 3.95, 8.8, 1.9,
    ["Dashboard    apollo-m.streamlit.app",
     "API                 apollo-api-tllm.onrender.com/docs",
     "Code               github.com/kashaffatimajaffrey-design/apollo-m",
     "CEREBRO       cerebro-sandy-beta.vercel.app"],
    size=15, align=PP_ALIGN.CENTER, color=BLUE, bold=True, space=9)

# ───────────────────────── Speaker notes ────────────────────────────────
# One short spoken version of each slide. Deliberately conversational and
# roughly the length of what fits comfortably in a minute, so it can be
# rehearsed rather than read.
NOTES = [
    "Good morning. We are presenting APOLLO-M, a framework that forecasts "
    "instability in online communities. I am Kashaf, with Rizwan and Aaqib, "
    "supervised by Mr. Saghir Ahmed. Everything you will see is deployed and "
    "publicly reachable, so you can open it yourself.",

    "I will start with the problem, then our idea, then what we built. The "
    "middle of the talk is our main result and the evidence behind it. At the "
    "end: who would use this, our timeline, what we cannot yet claim, and "
    "where we go next.",

    "Online communities do break down. Hostility rises, the good members "
    "leave, and moderators are left firefighting. The tools that exist judge "
    "one comment at a time, after it is posted. They describe the present. "
    "Nobody tells a moderator which community is heading downhill.",

    "Our starting point was Universe 25. Calhoun gave mice everything except "
    "space, and the colony collapsed. It was described afterwards; nothing "
    "measured it coming. Online communities have the same problem, except "
    "everything is recorded. So the warning signs are already in the data. "
    "We treat this as forecasting, not classification.",

    "Five layers. We score each comment, roll that up per community, group "
    "similar communities, forecast five days ahead, and raise a graded alert. "
    "The rule we followed throughout: code and trained models make every "
    "decision, and the language model only writes the explanation.",

    "In plain terms: micro scores comments with toxic-bert. Meso turns that "
    "into a health score. The pattern layer clusters communities with no "
    "labels. Macro forecasts with a Temporal Fusion Transformer, and gives "
    "confidence bands, not one number. Act turns it into an alert and a "
    "recommended action.",

    "This is our main result, and it is not the one we expected. We hid "
    "unstable communities in the data and asked each signal to rank them. "
    "Our own Community Health Index scored 0.575 — worse than raw toxicity, "
    "which is one of its own ingredients. The reason is that a health score "
    "says how bad things are now; instability is about how fast they are "
    "changing. Once we measured change, we got 1.000.",

    "Then we checked it against our graph model, which was built separately "
    "with different inputs. It scored 0.415, below chance, and for exactly "
    "the same reason. One weak score could be a bug in one formula. Two, "
    "arrived at independently, mean the whole type of measure is wrong for "
    "this task. That generalisation is what makes it a finding.",

    "Yes. We planted 15 destabilising communities among 60, and the "
    "forecaster never sees which. It got the direction right for all 15, and "
    "ranked them at 0.732. We tune for recall on purpose: missing a "
    "community that is deteriorating costs far more than checking a quiet "
    "one that turns out fine.",

    "This is not only simulation. We collected 116,854 genuine Reddit "
    "comments across 24 communities over 40 days, through Arctic Shift, with "
    "no API key. The whole pipeline runs on it. r/PublicFreakout came out "
    "least healthy — we never told the system that, it found it.",

    "All of this is live. There is a REST API with interactive documentation, "
    "Grafana monitoring reading real metrics, a button that fetches Reddit "
    "comments while you watch, and the five-day forecast with its uncertainty "
    "band. You can open any of these links right now.",

    "For evidence: 14 automated tests on the forecasting pipeline, and 67 on "
    "CEREBRO, the companion system. The toxicity classifier gets 90.1% on "
    "real labelled data we held back. And every number in our report is "
    "produced by a script, so it can be checked by re-running it.",

    "Who is this for? Trust and Safety teams first — they cannot watch 500 "
    "communities at once, and we tell them which ten matter this week. Then "
    "community platform vendors who can resell it, brand safety firms whose "
    "advertisers pay to avoid these spaces, and large organisations with "
    "internal communities. The hard part is knowing what to measure, and "
    "that is the part we have done.",

    "The industry is Trust and Safety. On the goals: SDG 16, because online "
    "hostility comes before real-world harm, and acting earlier means acting "
    "before it spreads. And SDG 9, because the API, the monitoring and the "
    "published benchmark are reusable infrastructure, not a one-off model.",

    "Eleven workstreams, all delivered. The honest part of this chart is "
    "weeks 4 to 6 — they were not in our plan. We audited our own data, "
    "found real problems, and stopped to fix them before continuing.",

    "What we will not claim. Our ground truth is planted, because nobody "
    "publishes a list of genuinely unstable communities — so 1.000 shows the "
    "pipeline is wired correctly, not real-world accuracy. The "
    "misinformation model was trained on news articles and does not transfer "
    "to short posts, so we publish no rate from it. Evaluation is "
    "single-fold, and English only.",

    "Next: real ground truth, using Reddit's own quarantine and ban dates, "
    "training only on what came before and checking whether we flagged it "
    "first. Longer history. Wiring the graph model in now that we know what "
    "it contributes. Retraining the two out-of-domain models. And publishing "
    "the benchmark so others can reproduce us.",

    "To close: we set out to predict breakdown rather than react to it. We "
    "built it and deployed it. We tested it honestly and found our own score "
    "was the weakest signal we had — then fixed it, and confirmed the same "
    "failure in a second model. It runs on real data. Anyone can build a "
    "pipeline; measuring yours and reporting what it says is the harder part.",

    "Thank you. All the links are here and everything is live. We are happy "
    "to take questions.",
]

for _s, _note in zip(prs.slides, NOTES):
    _s.notes_slide.notes_text_frame.text = _note

prs.save(str(OUT))
print(f"saved -> {OUT}")
print(f"slides: {len(prs.slides)} | notes: {len(NOTES)}")
